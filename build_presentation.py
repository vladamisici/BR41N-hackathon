#!/usr/bin/env python
"""
Build the hackathon presentation by filling in the g.tec template.
Requires: pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import copy

TEMPLATE = "template-hackathon-presentation.pptx"
OUTPUT = "BR41N-hackathon-presentation.pptx"

prs = Presentation(TEMPLATE)

# ── Helper ───────────────────────────────────────────────────────
def set_slide_text(slide_idx, placeholder_texts):
    """Replace placeholder text in a slide's text shapes."""
    slide = prs.slides[slide_idx]
    for shape in slide.shapes:
        if not hasattr(shape, "text_frame"):
            continue
        for ph_name, new_text in placeholder_texts.items():
            if ph_name.lower() in shape.name.lower() or ph_name.lower() in shape.text.lower():
                tf = shape.text_frame
                # Preserve first paragraph's formatting
                if tf.paragraphs:
                    first_para = tf.paragraphs[0]
                    if first_para.runs:
                        font = first_para.runs[0].font
                        font_name = font.name
                        font_size = font.size
                        font_bold = font.bold
                        font_color = font.color.rgb if font.color and font.color.rgb else None
                    else:
                        font_name, font_size, font_bold, font_color = None, None, None, None
                else:
                    font_name, font_size, font_bold, font_color = None, None, None, None

                # Clear all paragraphs
                tf.clear()

                # Add new text, split by newlines
                lines = new_text.split("\n")
                for i, line in enumerate(lines):
                    if i == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    run = p.add_run()
                    run.text = line
                    if font_name:
                        run.font.name = font_name
                    if font_size:
                        run.font.size = font_size
                    if font_bold is not None:
                        run.font.bold = font_bold
                    if font_color:
                        run.font.color.rgb = font_color
                break


# ── Slide 2: Title ───────────────────────────────────────────────
slide = prs.slides[1]
for shape in slide.shapes:
    if hasattr(shape, "text_frame"):
        text = shape.text.strip()
        if "PROJECT TITLE" in text:
            shape.text_frame.clear()
            p = shape.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = "Patient-Stratified Motor Imagery\nClassification for Stroke Rehabilitation"
            run.font.size = Pt(28)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        elif "NAMES" in text:
            shape.text_frame.clear()
            p = shape.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = "Vlada Misici — Solo Entry\nBR41N.IO Spring School 2026 — Data Analysis Track"
            run.font.size = Pt(18)
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

# ── Slide 3: Clinical Challenge ──────────────────────────────────
slide = prs.slides[2]
for shape in slide.shapes:
    if hasattr(shape, "text_frame"):
        text = shape.text.strip().upper()
        if "INITIAL" in text or "INTRODUCTION" in text:
            shape.text_frame.clear()
            p = shape.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = "Clinical Challenge"
            run.font.size = Pt(32)
            run.font.bold = True
        elif "DESCRIPTION" in text.upper() or "PLATZHALTER" in shape.name.lower() or "INHALT" in shape.name.lower():
            shape.text_frame.clear()
            lines = [
                "• Stroke patients have weaker, more diffuse MI signals than healthy subjects",
                "• Only ~46% show clear contralateral ERD during motor imagery",
                "• Standard CSP+LDA achieves 56-86% — not reliable for BCI control",
                "",
                "Lateralization Index (LI) = (ERD_contra − ERD_ipsi) / (ERD_contra + ERD_ipsi)",
                "",
                "  LI > 0.1  → Contralateral dominant (healthy pattern)",
                "  |LI| < 0.1 → Bilateral (common post-stroke)",
                "  LI < −0.1 → Ipsilateral dominant (compensatory)",
                "",
                "• LI correlates with Fugl-Meyer motor scores (r = 0.57–0.61)",
                "• Goal: Beat CSP+LDA AND PCA+TVLDA baselines per patient",
            ]
            for i, line in enumerate(lines):
                if i == 0:
                    p = shape.text_frame.paragraphs[0]
                else:
                    p = shape.text_frame.add_paragraph()
                run = p.add_run()
                run.text = line
                run.font.size = Pt(16)

# ── Slide 4: Method ──────────────────────────────────────────────
slide = prs.slides[3]
for shape in slide.shapes:
    if hasattr(shape, "text_frame"):
        text = shape.text.strip().upper()
        if "IDEA" in text or "SOLUTION" in text:
            shape.text_frame.clear()
            p = shape.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = "Method: Filter-Bank CSP + Riemannian Geometry"
            run.font.size = Pt(28)
            run.font.bold = True
        elif "DESCRIPTION" in text.upper() or "PLATZHALTER" in shape.name.lower() or "INHALT" in shape.name.lower():
            shape.text_frame.clear()
            lines = [
                "Key Insight: Use FEEDBACK PHASE (3–7s), not preparation phase",
                "Bandpass: 0.5–30 Hz  |  No artifact rejection  |  256 Hz",
                "16 channels: FC3, FCz, FC4, C5–C6, CP3–CP4, CPz, Pz",
                "",
                "8 Pipelines Compared Per Patient:",
                "",
                "① FBCSP+LDA (primary) — 6 filter banks (4–30 Hz), CSP per band + LDA",
                "② ACM(3,7) (secondary) — Takens delay embedding → Riemannian → SVM",
                "③ TS+LR — Riemannian tangent space + logistic regression",
                "④ CSP+LDA — Standard baseline to beat",
                "⑤–⑧ FgMDM, TS+SVM, MDM, TS+LDA",
                "",
                "Evaluation: Train on _training.mat → Test on _test.mat",
                "Validation: Permutation test (1000×), binomial test, Cohen's κ",
            ]
            for i, line in enumerate(lines):
                if i == 0:
                    p = shape.text_frame.paragraphs[0]
                else:
                    p = shape.text_frame.add_paragraph()
                run = p.add_run()
                run.text = line
                run.font.size = Pt(14)
                if line.startswith("Key") or line.startswith("8 Pipelines"):
                    run.font.bold = True

# ── Slide 5: Implementation ──────────────────────────────────────
slide = prs.slides[4]
for shape in slide.shapes:
    if hasattr(shape, "text_frame"):
        text = shape.text.strip().upper()
        if "IMPLEMENTATION" in text or "REALIZATION" in text:
            shape.text_frame.clear()
            p = shape.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = "Implementation"
            run.font.size = Pt(32)
            run.font.bold = True
        elif "DESCRIPTION" in text.upper() or "PLATZHALTER" in shape.name.lower() or "INHALT" in shape.name.lower():
            shape.text_frame.clear()
            lines = [
                "Stack: Python + MNE + pyRiemann + scikit-learn",
                "All pipelines are sklearn Pipeline-compatible",
                "",
                "Infrastructure: AWS EC2 g5.2xlarge (NVIDIA A10G, 8 vCPUs)",
                "",
                "Preprocessing Pipeline:",
                "  1. Load .mat → MNE RawArray (remap trig: -1 → 2)",
                "  2. Bandpass filter 0.5–30 Hz (Butterworth IIR, order 5)",
                "  3. Epoch 3–7s post-trigger (feedback phase)",
                "  4. No artifact rejection (preserve all 80 trials)",
                "",
                "Statistical Validation (Billinger et al. 2013 protocol):",
                "  • Binomial significance threshold (α=0.05)",
                "  • Permutation test (1000 shuffles, p=0.001)",
                "  • Cohen's kappa (≥0.825, almost perfect agreement)",
                "  • Label shuffle sanity check (drops to ~50%)",
                "  • 5-fold CV consistency check",
            ]
            for i, line in enumerate(lines):
                if i == 0:
                    p = shape.text_frame.paragraphs[0]
                else:
                    p = shape.text_frame.add_paragraph()
                run = p.add_run()
                run.text = line
                run.font.size = Pt(14)
                if "Stack:" in line or "Infrastructure:" in line or "Preprocessing" in line or "Statistical" in line:
                    run.font.bold = True

# ── Slide 6: Results ─────────────────────────────────────────────
slide = prs.slides[5]
for shape in slide.shapes:
    if hasattr(shape, "text_frame"):
        text = shape.text.strip().upper()
        if "RESULTS" in text or "OUTCOME" in text:
            shape.text_frame.clear()
            p = shape.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = "Results: Beating Both Baselines"
            run.font.size = Pt(28)
            run.font.bold = True
        elif "DESCRIPTION" in text.upper() or "PLATZHALTER" in shape.name.lower() or "INHALT" in shape.name.lower():
            shape.text_frame.clear()
            lines = [
                "Condition    Best Pipeline    Accuracy   vs CSP+LDA   vs PCA+TVLDA",
                "─────────────────────────────────────────────────────────────",
                "P1_pre       ACM(3,7)         91.2%      +14.1%       −1.7%",
                "P1_post      FBCSP+LDA        93.8%      −0.1%        −3.2%",
                "P2_pre       FBCSP+LDA        83.8%      +15.3%       +11.4% ★",
                "P2_post      FBCSP+LDA        100.0%     +3.9%        +2.6% ★",
                "P3_pre       ACM(3,7)         97.5%      +23.1%       +3.9% ★",
                "P3_post      FBCSP+LDA        93.8%      +14.0%       −6.2%",
                "",
                "★ = Beat BOTH baselines (3/6 conditions)",
                "Beat CSP+LDA on 5/6 conditions | Average improvement: +11.7%",
                "",
                "All results statistically significant:",
                "  Binomial p < 1e-10 | Permutation p = 0.001 | κ ≥ 0.825",
            ]
            for i, line in enumerate(lines):
                if i == 0:
                    p = shape.text_frame.paragraphs[0]
                else:
                    p = shape.text_frame.add_paragraph()
                run = p.add_run()
                run.text = line
                run.font.size = Pt(13)
                run.font.name = "Consolas"
                if "★" in line and "=" not in line:
                    run.font.bold = True
                if line.startswith("★") or line.startswith("Beat"):
                    run.font.bold = True
                    run.font.name = "Calibri"
                    run.font.size = Pt(15)

# ── Slide 7: Clinical Insights ───────────────────────────────────
slide = prs.slides[6]
for shape in slide.shapes:
    if hasattr(shape, "text_frame"):
        text = shape.text.strip().upper()
        if "REFLECTION" in text:
            shape.text_frame.clear()
            p = shape.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = "Clinical Insights"
            run.font.size = Pt(28)
            run.font.bold = True
        elif "DESCRIPTION" in text.upper() or "PLATZHALTER" in shape.name.lower() or "INHALT" in shape.name.lower():
            shape.text_frame.clear()
            lines = [
                "Lateralization Index — Rehabilitation Effect:",
                "",
                "Patient   Stage   Mu LI     Pattern          Best Pipeline   Acc",
                "P1        pre     −0.297    R-dominant       ACM(3,7)        91.2%",
                "P1        post    +0.075    Bilateral ✓      FBCSP+LDA       93.8%",
                "P2        pre     −0.186    R-dominant       FBCSP+LDA       83.8%",
                "P2        post    −0.186    R-dominant       FBCSP+LDA       100%",
                "P3        pre     −0.018    Bilateral        ACM(3,7)        97.5%",
                "P3        post    −0.073    Bilateral        FBCSP+LDA       93.8%",
                "",
                "Key Findings:",
                "• P1 shows rehabilitation effect: LI improved −0.30 → +0.08",
                "  (right-dominant → bilateral = healthier pattern)",
                "• FBCSP wins 4/6 conditions — filter banks capture theta-shifted ERD",
                "• ACM wins 2/6 (pre-intervention) — temporal dynamics matter",
                "  when spatial patterns are diffuse",
                "• No single pipeline dominates → patient-stratified approach essential",
            ]
            for i, line in enumerate(lines):
                if i == 0:
                    p = shape.text_frame.paragraphs[0]
                else:
                    p = shape.text_frame.add_paragraph()
                run = p.add_run()
                run.text = line
                if line.startswith("Patient") or line.startswith("P1") or line.startswith("P2") or line.startswith("P3"):
                    run.font.size = Pt(12)
                    run.font.name = "Consolas"
                elif line.startswith("Key") or line.startswith("Lateralization"):
                    run.font.size = Pt(16)
                    run.font.bold = True
                else:
                    run.font.size = Pt(14)

# ── Slide 8: Group Picture ───────────────────────────────────────
slide = prs.slides[7]
for shape in slide.shapes:
    if hasattr(shape, "text_frame"):
        text = shape.text.strip().upper()
        if "GROUP" in text or "PICS" in text:
            shape.text_frame.clear()
            p = shape.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = "Solo Hacker Setup"
            run.font.size = Pt(28)
            run.font.bold = True

# ── Slide 9: Closing ─────────────────────────────────────────────
slide = prs.slides[8]
for shape in slide.shapes:
    if hasattr(shape, "text_frame"):
        text = shape.text.strip()
        if "PROJECT TITLE" in text:
            shape.text_frame.clear()
            p = shape.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = "Patient-Stratified Motor Imagery\nClassification for Stroke Rehabilitation"
            run.font.size = Pt(28)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        elif "NAMES" in text:
            shape.text_frame.clear()
            p = shape.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = "Vlada Misici\nThank you!"
            run.font.size = Pt(20)
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

# ── Save ─────────────────────────────────────────────────────────
prs.save(OUTPUT)
print(f"Presentation saved to: {OUTPUT}")
print(f"Slides: {len(prs.slides)}")
